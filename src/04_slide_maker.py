#!/usr/bin/env python3
"""
Optimized PowerPoint Generator with Native Charts and Tables
"""

import argparse
import os
import sys
import logging
from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('slide_maker.log'),
        logging.StreamHandler()
    ]
)

def add_bar_chart(slide, df, x, y):
    chart_data = CategoryChartData()
    chart_data.categories = df.index.astype(str)
    chart_data.add_series('PV', df["PV"])
    cx, cy = Inches(9), Inches(4.5)
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)

def validate_inputs(args):
    """Check all input files exist and are valid"""
    if not os.path.exists(args.valuation):
        logging.error(f"Valuation file not found: {args.valuation}")
        sys.exit(1)
        
    if args.template and not os.path.exists(args.template):
        logging.warning(f"Template not found, using default: {args.template}")
        args.template = None

def get_default_template():
    """Create a basic default template"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Investment Thesis: SEA Acquisition of 11Street Korea"
    return prs

def build_presentation(template_path, valuation_path):
    """Create PowerPoint presentation with native charts and tables"""
    try:
        template_path = str(template_path) if template_path else None
        valuation_path = str(valuation_path)
        
        if template_path and os.path.exists(template_path):
            prs = Presentation(template_path)
        else:
            prs = get_default_template()
        
        # Read data
        dcf_data = pd.read_excel(valuation_path, sheet_name='DCF', index_col=0)
        merger_data = pd.read_excel(valuation_path, sheet_name='Merger_Model')
        
        # Valuation slide with bar chart
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "DCF Valuation Summary"
        
        plot_df = dcf_data[dcf_data.index != "Enterprise Value"].copy()
        plot_df['PV'] = plot_df['PV'] / 1e6  # Millions
        add_bar_chart(slide, plot_df, Inches(0.5), Inches(1.5))
        
        # Merger slide with table
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Accretion/Dilution Analysis"
        
        left = Inches(0.5)
        top = Inches(1.5)
        rows = len(merger_data) + 1
        cols = 2
        table = slide.shapes.add_table(rows, cols, left, top, Inches(9), Inches(3))
        table.first_row = True
        
        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Value"
        
        for i, row in merger_data.iterrows():
            table.cell(i+1, 0).text = row["Metric"]
            value = row["Value"]
            if "%" in row["Metric"]:
                table.cell(i+1, 1).text = f"{value:.1f}%"
            else:
                table.cell(i+1, 1).text = f"{value:.2f}"
        
        # TODO: Add more slides as needed to reach 10-slide pitchbook (e.g., comps, precedents, rationale)
        
        return prs
        
    except Exception as e:
        logging.error(f"Presentation build failed: {str(e)}")
        raise

def libreoffice_available():
    """Check if LibreOffice is installed"""
    import shutil
    return shutil.which("soffice") is not None
    
def convert_with_libreoffice(input_path, output_path):
    """Convert PPTX to PDF using LibreOffice"""
    try:
        import subprocess
        input_path = str(input_path)
        output_dir = str(Path(output_path).parent)
        
        result = subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf", 
             input_path, "--outdir", output_dir],
            capture_output=True,
            text=True,
            timeout=60
        )
        
        if result.returncode == 0:
            temp_pdf = Path(input_path).with_suffix('.pdf')
            if temp_pdf.exists():
                os.rename(str(temp_pdf), str(output_path))
                return True
        return False
        
    except Exception as e:
        logging.error(f"LibreOffice conversion error: {str(e)}")
        return False

def create_fallback_pdf(output_path):
    """Create simple fallback PDF if conversion fails"""
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
        
        c = canvas.Canvas(str(output_path), pagesize=letter)
        c.setFont("Helvetica-Bold", 16)
        c.drawString(100, 750, "Pitchbook Summary")
        c.setFont("Helvetica", 12)
        c.drawString(100, 700, "Please view the PPTX file for full details.")
        c.drawString(100, 650, "Valuation and merger models are in the Excel file.")
        c.save()
        
    except Exception as e:
        logging.error(f"Fallback PDF creation failed: {str(e)}")
        raise

def main():
    parser = argparse.ArgumentParser(description="Pitchbook Generator")
    parser.add_argument("--template", help="PPTX template file")
    parser.add_argument("--valuation", required=True, help="Valuation Excel file (includes merger)")
    parser.add_argument("--output", required=True, help="Output PDF file")
    args = parser.parse_args()
    
    try:
        validate_inputs(args)
        output_path = Path(args.output)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        pptx_path = output_path.with_suffix('.pptx')
        prs = build_presentation(args.template, args.valuation)
        prs.save(str(pptx_path))
        
        # Convert to PDF
        if libreoffice_available() and convert_with_libreoffice(pptx_path, output_path):
            logging.info("PDF successfully created via LibreOffice")
            os.remove(str(pptx_path))  # Clean up PPTX
        else:
            logging.warning("Using fallback PDF method")
            create_fallback_pdf(output_path)
                
        logging.info(f"Successfully created: {output_path}")
        
    except Exception as e:
        logging.error(f"Fatal error: {str(e)}")
        sys.exit(1)

if __name__ == "__main__":
    main()